"""Part 1: Helper functions for PPT builder"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Theme colors
BLACK = RGBColor(0x0a, 0x0a, 0x0a)
DARK = RGBColor(0x14, 0x14, 0x14)
CARD_BG = RGBColor(0x1a, 0x1a, 0x1a)
GOLD = RGBColor(0xd4, 0xaf, 0x37)
GOLD_LIGHT = RGBColor(0xf5, 0xd0, 0x60)
RED = RGBColor(0xe6, 0x39, 0x46)
TEAL = RGBColor(0x2e, 0xc4, 0xb6)
WHITE = RGBColor(0xe8, 0xe6, 0xe3)
GRAY = RGBColor(0xa0, 0x99, 0x8a)
DARK_GRAY = RGBColor(0x66, 0x66, 0x66)

def set_slide_bg(slide, color=BLACK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, font_size=14,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_card(slide, left, top, width, height, title, value, color=GOLD, bg=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = color
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.size = Pt(20)
    p2.font.color.rgb = color
    p2.font.bold = True
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER

def add_gold_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

def add_bullet_list(slide, left, top, width, height, items, font_size=12, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(4)

print("Part 1 loaded OK")
